"""Shared slide-building helpers for the DGE-METRIC Online Training decks.

GIZ color scheme matches ../Training Slides/def.tex (gizred RGB 204,7,30).
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

GIZ_RED = RGBColor(204, 7, 30)
GIZ_RED_DARK = RGBColor(150, 5, 22)
GIZ_GREEN = RGBColor(0, 104, 56)
LIGHT_GRAY = RGBColor(245, 245, 245)
MID_GRAY = RGBColor(120, 120, 120)
DARK = RGBColor(35, 35, 35)
WHITE = RGBColor(255, 255, 255)
AMBER = RGBColor(178, 106, 0)
AMBER_BG = RGBColor(255, 244, 229)

LOGO_PATH = os.path.normpath(os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "..", "..", "pictures", "GIZ_Logo.png"))

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT = "Calibri"


def new_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _set_font(run, size=18, color=DARK, bold=False, italic=False, font=FONT):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font


def _no_autosize(tf):
    tf.word_wrap = True
    try:
        tf.auto_size = None
    except Exception:
        pass


def add_footer(slide, page_no, deck_label):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, SLIDE_H - Inches(0.32), SLIDE_W, Inches(0.32))
    bar.fill.solid()
    bar.fill.fore_color.rgb = WHITE
    bar.line.fill.background()
    bar.shadow.inherit = False

    tb = slide.shapes.add_textbox(Inches(0.5), SLIDE_H - Inches(0.34), Inches(8), Inches(0.3))
    tf = tb.text_frame
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = deck_label
    _set_font(r, size=10, color=MID_GRAY)

    tb2 = slide.shapes.add_textbox(SLIDE_W - Inches(1.3), SLIDE_H - Inches(0.34), Inches(0.9), Inches(0.3))
    tf2 = tb2.text_frame
    tf2.margin_top = 0
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run()
    r2.text = str(page_no)
    _set_font(r2, size=10, color=MID_GRAY)


HEADER_H = Inches(1.4)

def add_header(slide, title, kicker=None):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, HEADER_H)
    band.fill.solid()
    band.fill.fore_color.rgb = GIZ_RED
    band.line.fill.background()
    band.shadow.inherit = False

    left = Inches(0.55)
    tb = slide.shapes.add_textbox(left, Inches(0.1), Inches(9.6), HEADER_H - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    if kicker:
        p0 = tf.paragraphs[0]
        r0 = p0.add_run()
        r0.text = kicker.upper()
        _set_font(r0, size=11, color=RGBColor(255, 210, 210), bold=True)
        p0.space_after = Pt(14)
        p1 = tf.add_paragraph()
        r1 = p1.add_run()
        r1.text = title
        _set_font(r1, size=24, color=WHITE, bold=True)
    else:
        p1 = tf.paragraphs[0]
        r1 = p1.add_run()
        r1.text = title
        _set_font(r1, size=26, color=WHITE, bold=True)

    logo = slide.shapes.add_picture(LOGO_PATH, SLIDE_W - Inches(1.9), (HEADER_H - Inches(0.48)) // 2, height=Inches(0.48))


def content_slide(prs, title, kicker=None, page_no=None, deck_label=""):
    slide = _blank(prs)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    bg.shadow.inherit = False
    add_header(slide, title, kicker)
    if page_no is not None:
        add_footer(slide, page_no, deck_label)
    return slide


BODY_TOP = HEADER_H + Inches(0.3)
BODY_LEFT = Inches(0.55)
BODY_W = SLIDE_W - Inches(1.1)
BODY_H = SLIDE_H - BODY_TOP - Inches(0.45)


def add_bullets(slide, items, top=BODY_TOP, left=BODY_LEFT, width=BODY_W, height=BODY_H,
                 size=17, color=DARK):
    """items: list of (text, level, bold, color_override) tuples, level in bullet-indent flavour."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        text, level, bold = item[0], item[1], item[2]
        col = item[3] if len(item) > 3 and item[3] else color
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = 0
        prefix = "•  " if level == 0 else ("–  " if level == 1 else "·  ")
        indent = Inches(0.0) if level == 0 else Inches(0.35 * level)
        p.space_after = Pt(10 if level == 0 else 6)
        r = p.add_run()
        r.text = prefix + text
        _set_font(r, size=size - (2 if level >= 1 else 0), color=col, bold=bold)
        pPr = p._pPr if p._pPr is not None else p.get_or_add_pPr()
        pPr.set('marL', str(int(indent)))
        pPr.set('indent', '0')
    return tb


def bullets_slide(prs, title, items, kicker=None, page_no=None, deck_label="", note=None):
    slide = content_slide(prs, title, kicker=kicker, page_no=page_no, deck_label=deck_label)
    h = BODY_H - (Inches(0.55) if note else Inches(0))
    add_bullets(slide, items, height=h)
    if note:
        add_note_box(slide, note, top=SLIDE_H - Inches(1.0))
    return slide


def add_note_box(slide, text, top, color=GIZ_RED, bg=LIGHT_GRAY, left=BODY_LEFT, width=BODY_W, height=Inches(0.62), size=13):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = bg
    box.line.color.rgb = color
    box.line.width = Pt(1)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    _set_font(r, size=size, color=DARK, italic=True)
    return box


def warning_slide(prs, title, headline, bullets, kicker=None, page_no=None, deck_label=""):
    """Prominent amber callout slide, used for the reserve-requirement scope caveat."""
    slide = content_slide(prs, title, kicker=kicker, page_no=page_no, deck_label=deck_label)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, BODY_LEFT, BODY_TOP, BODY_W, Inches(1.15))
    box.fill.solid()
    box.fill.fore_color.rgb = AMBER_BG
    box.line.color.rgb = AMBER
    box.line.width = Pt(1.5)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = headline
    _set_font(r, size=19, color=AMBER, bold=True, italic=True)
    add_bullets(slide, bullets, top=BODY_TOP + Inches(1.35), height=BODY_H - Inches(1.35))
    return slide


def table_slide(prs, title, headers, rows, kicker=None, page_no=None, deck_label="",
                 col_widths=None, font_size=13, header_color=GIZ_RED, note=None):
    slide = content_slide(prs, title, kicker=kicker, page_no=page_no, deck_label=deck_label)
    nrows = len(rows) + 1
    ncols = len(headers)
    top = BODY_TOP
    height = BODY_H - (Inches(0.55) if note else Inches(0))
    gtable = slide.shapes.add_table(nrows, ncols, BODY_LEFT, top, BODY_W, height).table
    if col_widths:
        total = sum(col_widths)
        for i, w in enumerate(col_widths):
            gtable.columns[i].width = int(BODY_W * (w / total))
    for j, htext in enumerate(headers):
        cell = gtable.cell(0, j)
        cell.text = ""
        tf = cell.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = htext
        _set_font(r, size=font_size, color=WHITE, bold=True)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_top = Pt(4)
        cell.margin_bottom = Pt(4)
    for i, row in enumerate(rows, start=1):
        shade = WHITE if i % 2 == 1 else LIGHT_GRAY
        for j, val in enumerate(row):
            cell = gtable.cell(i, j)
            cell.text = ""
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = str(val)
            _set_font(r, size=font_size, color=DARK, bold=(j == 0))
            cell.fill.solid()
            cell.fill.fore_color.rgb = shade
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
    if note:
        add_note_box(slide, note, top=SLIDE_H - Inches(1.0))
    return slide


def two_col_slide(prs, title, left_title, left_items, right_title, right_items,
                   kicker=None, page_no=None, deck_label=""):
    slide = content_slide(prs, title, kicker=kicker, page_no=page_no, deck_label=deck_label)
    col_w = (BODY_W - Inches(0.4)) / 2
    lt = slide.shapes.add_textbox(BODY_LEFT, BODY_TOP, col_w, Inches(0.4))
    r = lt.text_frame.paragraphs[0].add_run()
    r.text = left_title
    _set_font(r, size=16, color=GIZ_RED, bold=True)
    add_bullets(slide, left_items, top=BODY_TOP + Inches(0.45), left=BODY_LEFT, width=col_w,
                height=BODY_H - Inches(0.45), size=15)

    rx = BODY_LEFT + col_w + Inches(0.4)
    rt = slide.shapes.add_textbox(rx, BODY_TOP, col_w, Inches(0.4))
    r2 = rt.text_frame.paragraphs[0].add_run()
    r2.text = right_title
    _set_font(r2, size=16, color=GIZ_RED, bold=True)
    add_bullets(slide, right_items, top=BODY_TOP + Inches(0.45), left=rx, width=col_w,
                height=BODY_H - Inches(0.45), size=15)
    return slide


def title_slide(prs, title, subtitle, meta_lines, page_no=1, deck_label=""):
    slide = _blank(prs)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    bg.shadow.inherit = False

    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.28))
    band.fill.solid()
    band.fill.fore_color.rgb = GIZ_RED
    band.line.fill.background()
    band.shadow.inherit = False

    logo = slide.shapes.add_picture(LOGO_PATH, Inches(0.6), Inches(0.55), height=Inches(0.62))

    tb = slide.shapes.add_textbox(Inches(0.9), Inches(2.5), Inches(11.5), Inches(2.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    _set_font(r, size=40, color=DARK, bold=True)
    p.space_after = Pt(14)
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = subtitle
    _set_font(r2, size=22, color=GIZ_RED, bold=False)

    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.95), Inches(5.15), Inches(4.2), Pt(3))
    rule.fill.solid()
    rule.fill.fore_color.rgb = GIZ_RED
    rule.line.fill.background()
    rule.shadow.inherit = False

    mtb = slide.shapes.add_textbox(Inches(0.9), Inches(5.35), Inches(11), Inches(1.6))
    mtf = mtb.text_frame
    mtf.word_wrap = True
    first = True
    for line in meta_lines:
        p = mtf.paragraphs[0] if first else mtf.add_paragraph()
        first = False
        r = p.add_run()
        r.text = line
        _set_font(r, size=15, color=MID_GRAY)
        p.space_after = Pt(4)

    bottomband = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, SLIDE_H - Inches(0.28), SLIDE_W, Inches(0.28))
    bottomband.fill.solid()
    bottomband.fill.fore_color.rgb = GIZ_RED
    bottomband.line.fill.background()
    bottomband.shadow.inherit = False
    return slide


def section_slide(prs, kicker, title, subtitle=None, page_no=None, deck_label=""):
    slide = _blank(prs)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = GIZ_RED
    bg.line.fill.background()
    bg.shadow.inherit = False

    logo = slide.shapes.add_picture(LOGO_PATH, SLIDE_W - Inches(2.3), Inches(0.5), height=Inches(0.5))

    tb = slide.shapes.add_textbox(Inches(0.9), Inches(2.9), Inches(11), Inches(2.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = kicker.upper()
    _set_font(r, size=15, color=RGBColor(255, 210, 210), bold=True)
    p.space_after = Pt(10)
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = title
    _set_font(r2, size=34, color=WHITE, bold=True)
    if subtitle:
        p3 = tf.add_paragraph()
        p3.space_before = Pt(14)
        r3 = p3.add_run()
        r3.text = subtitle
        _set_font(r3, size=17, color=WHITE)
    return slide
